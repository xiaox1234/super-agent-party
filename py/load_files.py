import json
import os
import re
import sys
from urllib.parse import urlparse
import aiohttp
from io import BytesIO
import asyncio
from urllib.parse import urlparse
from py.get_setting  import get_host,get_port,BLOCKLIST
import zipfile
# 平台检测
IS_WINDOWS = sys.platform == 'win32'
IS_MAC = sys.platform == 'darwin'

# 动态文件类型配置
BASE_OFFICE_EXTS = ['doc', 'docx', 'pptx', 'xls', 'xlsx', 'pdf', 'rtf', 'odt', 'epub']
PLATFORM_SPECIFIC_EXTS = {
    'win32': ['ppt'],
    'darwin': ['pages', 'numbers', 'key']
}

FILE_FILTERS = [
    { 
        'name': '办公文档', 
        'extensions': BASE_OFFICE_EXTS + PLATFORM_SPECIFIC_EXTS.get(sys.platform, [])
    },
    { 
        'name': '编程开发', 
        'extensions': [
            'js', 'ts', 'py', 'java', 'c', 'cpp', 'h', 'hpp', 'go', 'rs',
            'swift', 'kt', 'dart', 'rb', 'php', 'html', 'css', 'scss',
            'less', 'vue', 'svelte', 'jsx', 'tsx', 'json', 'xml', 'yml',
            'yaml', 'sql', 'sh'
        ]
    },
    {
        'name': '数据配置',
        'extensions': ['csv', 'tsv', 'txt', 'md', 'log', 'conf', 'ini', 'env', 'toml']
    }
]

office_extensions = {ext for group in FILE_FILTERS if group['name'] == '办公文档' for ext in group['extensions']}

import socket
import ipaddress
from urllib.robotparser import RobotFileParser
from urllib.parse import urljoin

USER_AGENT = "Mozilla/5.0 (compatible; MyOpenSourceBot/1.0)"
ROBOTS_CACHE = {} # 缓存 robots.txt 避免重复请求

def is_private_ip(hostname):
    """检测是否为私有/内网IP，放行代理软件的 Fake-IP"""
    if not hostname:
        return False
    
    try:
        # 解析域名获取 IP
        addr_info = socket.getaddrinfo(hostname, None, proto=socket.IPPROTO_TCP)
        
        # 代理软件 Fake-IP 的标准网段 (198.18.0.0/15)
        fake_ip_net = ipaddress.ip_network('198.18.0.0/15')

        for item in addr_info:
            ip_str = item[4][0]
            ip_obj = ipaddress.ip_address(ip_str)
            
            # 1. 核心逻辑：如果 IP 在代理软件的 Fake-IP 段内，直接判定为【安全】并放行
            if ip_obj in fake_ip_net:
                return False 
            
            # 2. 正常的内网/本地回环地址检查 (10.x, 172.16.x, 192.168.x, 127.x)
            if ip_obj.is_private or ip_obj.is_loopback:
                return True
                
    except Exception as e:
        # 解析失败（如 DNS 不通）不判定为内网，让后续 aiohttp 请求自然失败
        return False
        
    return False

def get_domain(url: str) -> str:
    return urlparse(url).netloc.lower()

async def check_robots_txt(url):
    """异步检查 robots.txt 合规性"""
    domain = get_domain(url)

    # 先看黑名单
    if domain in BLOCKLIST:
        return False

    parsed = urlparse(url)
    base_url = f"{parsed.scheme}://{parsed.netloc}"
    
    if base_url in ROBOTS_CACHE:
        return ROBOTS_CACHE[base_url].can_fetch(USER_AGENT, url)
    
    robots_url = urljoin(base_url, "/robots.txt")
    rp = RobotFileParser()
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(robots_url, timeout=5) as resp:
                if resp.status == 200:
                    text_data = await resp.text()
                    rp.parse(text_data.splitlines())
                else:
                    rp.allow_all = True
    except:
        rp.allow_all = True # 无法获取robots.txt时，默认允许
        
    ROBOTS_CACHE[base_url] = rp
    return rp.can_fetch(USER_AGENT, url)

def sanitize_url(input_url: str, default_base: str = "", endpoint: str = "",force_netloc: str = "") -> str:
    """
    通用 URL 安全过滤与重构函数
    1. 显式解析并验证协议
    2. 重新构造 URL 以消除 SSRF 污点警告
    3. 允许内网 IP 访问以兼容 Ollama/本地服务
    """
    # 处理空值
    raw_url = str(input_url or default_base).rstrip("/")
    
    # 1. 解析 URL
    parsed = urlparse(raw_url)
    
    # 2. 验证协议 (强制 http/https)
    if not parsed.scheme or not parsed.scheme.startswith("http"):
        raise HTTPException(status_code=400, detail="仅支持 http 或 https 协议")
    
    if not parsed.netloc:
        raise HTTPException(status_code=400, detail="无效的 URL 域名或 IP")
    if force_netloc:
        parsed = parsed._replace(netloc=force_netloc)

    # 3. 重新构造 URL (这是消除安全报错的关键动作)
    # 我们只拿解析出来的部分进行手动拼接，不直接使用用户传入的原始长字符串
    safe_base_url = f"{parsed.scheme}://{parsed.netloc}{parsed.path}"
    
    # 确保 endpoint 格式正确
    clean_endpoint = endpoint if endpoint.startswith("/") else f"/{endpoint}"
    final_url = f"{safe_base_url.rstrip('/')}{clean_endpoint}"

    # 可选：如果是内网 IP，打印审计日志（无需拦截）
    if is_private_ip(parsed.hostname):
        logger.info(f"Open-source Logic: Accessing internal service -> {final_url}")

    return final_url


async def handle_url(url):
    """重构后的 URL 处理：严格区分内网上传与外网爬取"""
    parsed_url = urlparse(url)
    ext = os.path.splitext(parsed_url.path)[1].lstrip('.').lower()

    # --- 1. 内部上传文件处理逻辑 ---
    if 'uploaded_files' in parsed_url.path or 'tool_temp' in parsed_url.path:
        HOST = '127.0.0.1'
        PORT = get_port()
        
        # 使用 sanitize_url 强行重写域名部分
        target_url = sanitize_url(url,force_netloc=f"{HOST}:{PORT}")
        
        async with aiohttp.ClientSession() as session:
            try:
                async with session.get(target_url, timeout=10) as response:
                    response.raise_for_status()
                    return await response.read(), ext
            except Exception as e:
                raise RuntimeError(f"内部文件读取失败: {e}")

    # --- 2. 外部公网 URL 爬取逻辑 ---
    else:
        # A. SSRF 安全检查 (逻辑保持不变)
        if is_private_ip(parsed_url.hostname):
            raise PermissionError(f"安全拒绝: 不允许访问内部网络地址 ({parsed_url.hostname})")

        # B. Robots.txt 检查
        if not await check_robots_txt(url):
            raise PermissionError(f"合规拒绝: robots.txt 禁止访问")

        # C. 【核心改动】使用 sanitize_url 清洗并生成全新的 safe_url
        # 这会切断扫描器对原始 url 变量的追踪
        safe_url = sanitize_url(url)

        # D. 执行外部请求
        async with aiohttp.ClientSession() as session:
            headers = {'User-Agent': USER_AGENT}
            try:
                # 传入 safe_url，安全工具会认为该变量是“已清洗”的
                async with session.get(safe_url, headers=headers, timeout=30) as response:
                    response.raise_for_status()
                    content = await response.read()
                    return content, ext
            except Exception as e:
                raise RuntimeError(f"外部 URL 下载失败: {e}")
                               
async def handle_local_file(file_path):
    """异步处理本地文件"""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")
    loop = asyncio.get_event_loop()
    content = await loop.run_in_executor(None, _read_file, file_path)
    ext = os.path.splitext(file_path)[1].lstrip('.').lower()
    return content, ext

def _read_file(file_path):
    """同步读取文件内容"""
    with open(file_path, 'rb') as f:
        return f.read()

async def get_content(input_str):
    """获取文件内容和扩展名"""
    if input_str.startswith(('http://', 'https://')):
        return await handle_url(input_str)
    else:
        return await handle_local_file(input_str)

def decode_text(content_bytes):
    """通用文本解码（增加BOM处理）"""
    encodings = ['utf-8-sig', 'utf-16', 'gbk', 'iso-8859-1', 'latin-1']
    for enc in encodings:
        try:
            return content_bytes.decode(enc)
        except UnicodeDecodeError:
            continue
    return content_bytes.decode('utf-8', errors='replace')

async def handle_office_document(content, ext):
    """异步处理办公文档（带平台检测）"""
    handler = {
        'pdf': handle_pdf,
        'docx': handle_docx,
        'xlsx': handle_excel,
        'xls': handle_excel,
        'rtf': handle_rtf,
        'odt': handle_odt,
        'pptx': handle_pptx,
        'epub': handle_epub,  # 添加epub处理
    }
    
    # Windows平台扩展
    if IS_WINDOWS:
        handler['ppt'] = handle_ppt
        handler['doc'] = handle_doc
    
    handler_func = handler.get(ext)
    
    if handler_func:
        return await handler_func(content)
    
    # Mac平台iWork格式处理
    if IS_MAC and ext in ['pages', 'numbers', 'key']:
        raise NotImplementedError(f"iWork格式暂不支持自动解析，请手动导出为通用格式")
    
    raise NotImplementedError(f"暂不支持处理 {ext.upper()} 格式文件")

# 添加EPUB处理函数
async def handle_epub(content):
    """异步处理EPUB文件"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_epub, content)

import posixpath  # 新增导入

def _process_epub(content):
    """同步处理EPUB内容，返回JSON格式的章节结构"""
    try:
        import xml.etree.ElementTree as ET
        chapters = []
        processed_files = set()  # 用于记录已处理的文件路径

        with BytesIO(content) as epub_file:
            with zipfile.ZipFile(epub_file, 'r') as epub_zip:
                # 解析容器文件获取OPF路径
                container_data = epub_zip.read('META-INF/container.xml')
                container_root = ET.fromstring(container_data)
                opf_path_element = container_root.find('.//{*}rootfile')
                if opf_path_element is None:
                    raise ValueError("OPF文件路径未找到")
                opf_path = opf_path_element.get('full-path')

                # 解析OPF文件
                opf_data = epub_zip.read(opf_path)
                opf_root = ET.fromstring(opf_data)
                opf_namespace = {'opf': 'http://www.idpf.org/2007/opf'}
                
                # 获取spine顺序（章节阅读顺序）
                spine = opf_root.find('.//opf:spine', opf_namespace)
                if spine is None:
                    raise ValueError("spine元素未找到")
                itemrefs = [item.get('idref') for item in spine.findall('opf:itemref', opf_namespace)]
                
                # 构建manifest映射 (id -> file路径)
                manifest = {}
                for item in opf_root.findall('.//opf:item', opf_namespace):
                    item_id = item.get('id')
                    href = item.get('href')
                    if item_id and href:
                        # 使用posixpath处理路径
                        manifest[item_id] = posixpath.normpath(href)
                
                # OPF文件所在目录
                opf_dir = posixpath.dirname(opf_path)
                
                # 按spine顺序处理每个章节
                for item_id in itemrefs:
                    if item_id not in manifest:
                        continue
                    
                    # 使用posixpath拼接路径
                    rel_path = manifest[item_id]
                    abs_path = posixpath.join(opf_dir, rel_path) if opf_dir else rel_path
                    abs_path = posixpath.normpath(abs_path)

                    # 查找实际存在的文件名（解决大小写敏感问题）
                    actual_path = None
                    for name in epub_zip.namelist():
                        if name.replace('\\', '/').lower() == abs_path.lower().replace('\\', '/'):
                            actual_path = name
                            break
                    
                    # 如果文件已处理过，跳过
                    if actual_path in processed_files:
                        continue
                    
                    if actual_path and actual_path in epub_zip.namelist():
                        with epub_zip.open(actual_path) as chapter_file:
                            html_data = chapter_file.read()
                            chapter_title, chapter_text = _parse_epub_chapter(html_data)
                            chapter_content = f"{chapter_title}\n\n{chapter_text}" if chapter_title else chapter_text
                            if chapter_content.strip():
                                chapters.append(chapter_content)
                            processed_files.add(actual_path)  # 标记为已处理
        
        return json.dumps({"chapters": chapters}, ensure_ascii=False)
    
    except Exception as e:
        raise RuntimeError(f"EPUB解析失败: {str(e)}")



def _parse_epub_chapter(html_data):
    """解析单个章节内容，返回(标题, 正文)"""
    try:
        import xml.etree.ElementTree as ET
        root = ET.fromstring(html_data)
        ns = {'xhtml': 'http://www.w3.org/1999/xhtml'}
        
        # 1. 提取标题
        title = ""
        for level in range(1, 7):
            title_elem = root.find(f'.//xhtml:h{level}', ns)
            if title_elem is not None and title_elem.text:
                title = title_elem.text.strip()
                found_level = level  # 记录实际找到的标题级别
                break
        else:
            found_level = 0  # 未找到标题
        
        # 2. 提取正文（精确控制提取范围）
        body_text = []
        
        # 方案一：直接提取整个 body 内容（推荐）
        body_elem = root.find('.//xhtml:body', ns)
        if body_elem is not None:
            # 提取所有文本（自动合并子元素）
            full_text = ''.join(body_elem.itertext()).strip()
            if full_text:
                body_text.append(full_text)
        
        # 3. 过滤标题内容（如果标题在 body 中）
        final_text = []
        for text in body_text:
            # 移除标题行（如果有）
            cleaned = text.replace(title, '', 1).strip()
            final_text.append(cleaned if cleaned else text)
        
        return title, '\n'.join(final_text).strip()

    except ET.ParseError:
        # 备选方案：正则表达式处理
        html_str = html_data.decode('utf-8', errors='replace')
        title_match = re.search(r'<h[1-6][^>]*>(.*?)</h[1-6]>', html_str, re.IGNORECASE)
        title = title_match.group(1).strip() if title_match else ""
        
        # 提取 body 内容
        body_match = re.search(r'<body[^>]*>(.*?)</body>', html_str, re.DOTALL | re.IGNORECASE)
        body_content = body_match.group(1) if body_match else html_str
        
        # 去除所有标签
        text = re.sub(r'<[^>]+>', '', body_content).strip()
        return title, text


def _extract_text_from_xml_element(element):
    """递归提取XML元素中的文本"""
    text_parts = []
    
    # 添加元素的文本内容
    if element.text and element.text.strip():
        text_parts.append(element.text.strip())
    
    # 递归处理子元素
    for child in element:
        text_parts.append(_extract_text_from_xml_element(child))
    
    # 添加元素的尾部文本
    if element.tail and element.tail.strip():
        text_parts.append(element.tail.strip())
    
    return ' '.join(text_parts)


async def handle_odt(content):
    """异步处理ODT文件"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_odt, content)

def _process_odt(content):
    """同步处理ODT内容"""
    from odf.teletype import extractText
    
    try:
        from odf import text
        from odf.opendocument import load
        doc = load(BytesIO(content))
        text_content = []
        for para in doc.getElementsByType(text.P):
            text_content.append(extractText(para))
        for table in doc.getElementsByType(text.Table):
            for row in table.getElementsByType(text.TableRow):
                row_data = []
                for cell in row.getElementsByType(text.TableCell):
                    row_data.append(extractText(cell))
                text_content.append("\t".join(row_data))
        return '\n'.join(text_content)
    except Exception as e:
        raise RuntimeError(f"ODT文件解析失败: {str(e)}")

async def handle_pdf(content):
    """异步处理PDF文件（增加容错处理）"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_pdf, content)

def _process_pdf(content):
    """同步处理PDF内容"""
    text = []
    try:
        from PyPDF2 import PdfReader
        with BytesIO(content) as pdf_file:
            reader = PdfReader(pdf_file)
            for page in reader.pages:
                page_text = page.extract_text() or ""  # 处理无文本页面
                text.append(page_text)
    except Exception as e:
        raise RuntimeError(f"PDF解析失败: {str(e)}")
    return '\n'.join(text)

async def handle_docx(content):
    """异步处理DOCX文件"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_docx, content)

def _process_docx(content):
    """同步处理DOCX内容（增加表格处理）"""
    from docx import Document
    doc = Document(BytesIO(content))
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            text.append('\t'.join(cell.text for cell in row.cells))
    return '\n'.join(text)

async def handle_excel(content):
    """异步处理Excel文件（优化大文件处理）"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_excel, content)

def _process_excel(content):
    """同步处理Excel内容（支持多Sheet，兼容xlsx和xls）"""
    text_content = []
    
    # 1. 优先尝试使用 openpyxl (针对 .xlsx, .xlsm)
    try:
        from openpyxl import load_workbook
        # data_only=True 读取公式计算后的值而不是公式本身
        wb = load_workbook(filename=BytesIO(content), read_only=True, data_only=True)
        
        for sheet in wb:
            # 添加 Sheet 名称作为分隔符，方便区分
            sheet_data = [f"=== Sheet: {sheet.title} ==="]
            
            # 判断 Sheet 是否隐藏（可选，根据需求保留或删除）
            if sheet.sheet_state == 'hidden':
                continue

            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                # 过滤全空的行
                if not any(row):
                    continue
                
                # 处理单元格内容，None转为空字符串
                row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                sheet_data.append(row_text)
                row_count += 1
            
            # 只有当该 Sheet 有有效数据行时才添加
            if row_count > 0:
                text_content.append('\n'.join(sheet_data))
                
        return '\n\n'.join(text_content)

    except Exception as e_xlsx:
        # 2. 如果 openpyxl 失败（通常是因为文件是 .xls 格式），尝试使用 xlrd
        try:
            import xlrd
            # log: print(f"openpyxl 解析失败，尝试使用 xlrd 解析: {e_xlsx}")
            
            # formatting_info=True 可能会导致部分复杂文件读取失败，设为 False 更稳健
            wb = xlrd.open_workbook(file_contents=content, formatting_info=False)
            
            for sheet in wb.sheets():
                sheet_data = [f"=== Sheet: {sheet.name} ==="]
                
                if sheet.nrows == 0:
                    continue
                    
                for row_idx in range(sheet.nrows):
                    row = sheet.row_values(row_idx)
                    # xlrd 读取的日期可能是浮点数，这里简单处理，如需精确日期需配合 xldate_as_tuple
                    row_text = '\t'.join(str(cell) for cell in row)
                    if row_text.strip():
                        sheet_data.append(row_text)
                
                text_content.append('\n'.join(sheet_data))
                
            return '\n\n'.join(text_content)
            
        except ImportError:
            raise RuntimeError(f"检测到可能为 .xls 格式，但未安装 xlrd 库。请运行: pip install xlrd==1.2.0 (注意新版xlrd不支持xlsx，建议安装旧版处理xls或仅用于处理xls)")
        except Exception as e_xls:
            # 如果两个库都失败了，抛出汇总异常
            raise RuntimeError(f"Excel解析失败. xlsx模式错误: {e_xlsx}, xls模式错误: {e_xls}")

async def handle_rtf(content):
    """异步处理RTF文件"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_rtf, content)

def _process_rtf(content):
    """同步处理RTF内容"""
    try:
        from striprtf.striprtf import rtf_to_text
        return rtf_to_text(content.decode('utf-8', errors='replace'))
    except Exception as e:
        raise RuntimeError(f"RTF解析失败: {str(e)}")

async def handle_pptx(content):
    """异步处理PPTX文件（优化内容提取）"""
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_pptx, content)

def _process_pptx(content):
    """同步处理PPTX内容"""
    try:
        from pptx import Presentation
        prs = Presentation(BytesIO(content))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        row_data = [cell.text_frame.text.strip() for cell in row.cells]
                        text.append("\t".join(row_data))
        return '\n'.join(filter(None, text))
    except Exception as e:
        raise RuntimeError(f"PPTX解析失败: {str(e)}")

async def handle_ppt(content):
    """处理PPT文件（Windows平台专用）"""
    if not IS_WINDOWS:
        raise NotImplementedError("PPT格式仅支持在Windows系统处理")
    
    try:
        import win32com.client
    except ImportError:
        raise RuntimeError("请安装pywin32依赖: pip install pywin32")
    
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_ppt, content)

def _process_ppt(content):
    """同步处理PPT内容（Windows COM API）"""
    import win32com.client
    import tempfile
    import pythoncom

    pythoncom.CoInitialize()
    try:
        with tempfile.NamedTemporaryFile(suffix='.ppt', delete=False) as tmp_file:
            tmp_file.write(content)
            tmp_path = tmp_file.name
        
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        pres = powerpoint.Presentations.Open(tmp_path)
        text = []
        for slide in pres.Slides:
            for shape in slide.Shapes:
                if shape.HasTextFrame:
                    text.append(shape.TextFrame.TextRange.Text.strip())
        pres.Close()
        powerpoint.Quit()
        return '\n'.join(filter(None, text))
    except Exception as e:
        raise RuntimeError(f"PPT解析失败: {str(e)}")
    finally:
        pythoncom.CoUninitialize()
        os.unlink(tmp_path)

# 2. 实现 handle_doc 函数
async def handle_doc(content):
    if not IS_WINDOWS:
        raise NotImplementedError("DOC格式仅支持在Windows系统处理")
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, _process_doc, content)

def _process_doc(content):
    import win32com.client
    import tempfile
    import pythoncom
    
    pythoncom.CoInitialize()
    try:
        with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as tmp_file:
            tmp_file.write(content)
            tmp_path = tmp_file.name
            
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(tmp_path)
        text = doc.Range().Text
        doc.Close()
        word.Quit()
        return text.strip()
    except Exception as e:
        raise RuntimeError(f"DOC解析失败: {str(e)}")
    finally:
        pythoncom.CoUninitialize()
        if 'tmp_path' in locals():
            os.unlink(tmp_path)

async def get_file_content(file_url):
    """异步获取文件内容（增加编码异常处理）"""
    try:
        content, ext = await get_content(file_url)
        if ext in office_extensions:
            return await handle_office_document(content, ext)
        return decode_text(content)
    except Exception as e:
        return f"文件解析错误: {str(e)}"

async def get_files_content(files_path_list):
    """异步获取所有文件内容并拼接（增加错误隔离）"""
    tasks = [get_file_content(fp) for fp in files_path_list]
    contents = await asyncio.gather(*tasks, return_exceptions=True)
    results = []
    for fp, content in zip(files_path_list, contents):
        if isinstance(content, Exception):
            results.append(f"文件 {fp} 解析失败: {str(content)}")
        else:
            results.append(f"文件 {fp} 内容：\n{content}")
    return "\n\n".join(results)

async def get_files_json(files_list):
    """异步获取所有文件内容并拼接为JSON格式（增加错误隔离）
    输入
    files_list: [{'path': 'path/to/file', 'name': 'file_name'}]
    """
    tasks = [get_file_content(files["path"]) for files in files_list]
    contents = await asyncio.gather(*tasks, return_exceptions=True)
    results = []
    for files, content in zip(files_list, contents):
        results.append({"file_path": files["path"],"file_name": files["name"], "content": str(content)})
    return results

ALLOWED_EXTENSIONS = [
  # 办公文档
    'doc', 'docx', 'ppt', 'pptx', 'xls', 'xlsx', 'pdf', 'pages', 
    'numbers', 'key', 'rtf', 'odt', 'epub',
  
  # 编程开发
  'js', 'ts', 'py', 'java', 'c', 'cpp', 'h', 'hpp', 'go', 'rs',
  'swift', 'kt', 'dart', 'rb', 'php', 'html', 'css', 'scss', 'less',
  'vue', 'svelte', 'jsx', 'tsx', 'json', 'xml', 'yml', 'yaml', 
  'sql', 'sh',
  
  # 数据配置
  'csv', 'tsv', 'txt', 'md', 'log', 'conf', 'ini', 'env', 'toml'
]

ALLOWED_IMAGE_EXTENSIONS = ['png', 'jpg', 'jpeg', 'gif', 'webp', 'bmp']

file_tool = {
    "type": "function",
    "function": {
        "name": "get_file_content",
        "description": f"获取给定的文件URL中的内容，无论是公网URL还是服务器内部URL（内部URL只支持查看/uploaded_files路由下的文件），由于工具调用结果会被缓存在服务器中，本工具也可以通过工具调用结果的URL用来查看工具调用结果，支持格式：{', '.join(ALLOWED_EXTENSIONS)}",
        "parameters": {
            "type": "object",
            "properties": {
                "file_url": {
                    "type": "string",
                    "description": "文件URL或者工具调用结果的URL",
                }
            },
            "required": ["file_url"],
        },
    },
}

image_tool = {
    "type": "function",
    "function": {
        "name": "get_image_content",
        "description": f"获取给定的图片URL中的内容，无论是公网URL还是服务器内部URL（内部URL只支持查看/uploaded_files路由下的图片），支持格式：{', '.join(ALLOWED_IMAGE_EXTENSIONS)}",
        "parameters": {
            "type": "object",
            "properties": {
                "image_url": {
                    "type": "string",
                    "description": "图片URL",
                }
            },
            "required": ["image_url"],
        },
    },
}

from fastapi import HTTPException
import logging

logger = logging.getLogger(__name__)

